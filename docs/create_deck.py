#!/usr/bin/env python3
"""
Generate the Superloop RWA Strategy Deck (PPTX)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ──────────────────────────────────────────────
BLACK      = RGBColor(0x0D, 0x0D, 0x0D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x00, 0x7A, 0xFF)   # Superloop blue
ACCENT2    = RGBColor(0x00, 0xC2, 0x8E)   # Teal/green
DARK_BG    = RGBColor(0x0F, 0x11, 0x1A)
DARK_CARD  = RGBColor(0x1A, 0x1D, 0x2E)
LIGHT_GREY = RGBColor(0xA0, 0xA0, 0xA8)
MED_GREY   = RGBColor(0x60, 0x60, 0x68)
HIGHLIGHT  = RGBColor(0xFF, 0xC1, 0x07)   # Gold / attention

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ───────────────────────────────────────────────────
def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, bold=False, spacing=1.2, bullet=False, alignment=PP_ALIGN.LEFT):
    """lines is list of (text, optional_color, optional_bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, clr, bld = item, color, bold
        elif len(item) == 2:
            txt, clr = item; bld = bold
        else:
            txt, clr, bld = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "  " if bullet else ""
        p.text = prefix + txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * (spacing - 1) * 2)
        p.alignment = alignment
    return txBox

def add_kpi_card(slide, left, top, width, height, title, value, subtitle=""):
    card = add_rounded_rect(slide, left, top, width, height, DARK_CARD, ACCENT)
    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3),
             title, font_size=11, color=LIGHT_GREY, bold=False)
    add_text(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.5),
             value, font_size=28, color=ACCENT, bold=True)
    if subtitle:
        add_text(slide, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), Inches(0.3),
                 subtitle, font_size=10, color=LIGHT_GREY)

def add_slide_number(slide, num, total):
    add_text(slide, W - Inches(1.2), H - Inches(0.45), Inches(1), Inches(0.35),
             f"{num} / {total}", font_size=9, color=MED_GREY, alignment=PP_ALIGN.RIGHT)

def add_header_bar(slide):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

TOTAL_SLIDES = 20

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "SUPERLOOP", font_size=56, color=ACCENT, bold=True)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
         "Institutional RWA Looping Infrastructure", font_size=32, color=WHITE, bold=False)
add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
         "Automated Leveraged Yield Vaults for Tokenized Real-World Assets",
         font_size=18, color=LIGHT_GREY)

add_rect(slide, Inches(1), Inches(4.3), Inches(4), Inches(0.03), ACCENT)

add_multiline(slide, Inches(1), Inches(4.8), Inches(6), Inches(1.5), [
    ("Built on Aave V3 | Aave Horizon | Centrifuge | Morpho", LIGHT_GREY),
    ("ERC-4626 Compliant | Modular Architecture | Flash-Loan Optimized", LIGHT_GREY),
    ("February 2026 | Confidential", MED_GREY),
], font_size=14)
add_slide_number(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "THE PROBLEM", font_size=32, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         "Tokenized RWAs are on-chain, but unlocking their full yield potential is manual, fragmented, and inaccessible.",
         font_size=16, color=LIGHT_GREY)

# Three problem cards
problems = [
    ("Manual Looping is Complex",
     "Institutions need 30+ transactions to loop RWA positions for leveraged yield. Each loop requires: swap, supply, borrow, repeat. Error-prone, gas-intensive, time-consuming."),
    ("Fragmented Infrastructure",
     "RWA tokens, lending markets, oracles, and compliance layers each live in separate protocols. No unified system to compose them into managed vault strategies."),
    ("Compliance Barriers",
     "RWA tokens carry transfer restrictions (ERC-1404, ERC-3643). Vaults must be KYC-whitelisted by each issuer. No plug-and-play compliance framework exists for DeFi vaults."),
]

for i, (title, desc) in enumerate(problems):
    left = Inches(0.8 + i * 3.9)
    card = add_rounded_rect(slide, left, Inches(1.8), Inches(3.6), Inches(3.5), DARK_CARD, MED_GREY)
    add_text(slide, left + Inches(0.25), Inches(2.0), Inches(3.1), Inches(0.5),
             title, font_size=18, color=HIGHLIGHT, bold=True)
    add_text(slide, left + Inches(0.25), Inches(2.7), Inches(3.1), Inches(2.3),
             desc, font_size=13, color=LIGHT_GREY)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
         "Result: $35B+ in tokenized RWAs sit under-utilized. Institutions earn base yield (~4-5%) when they could earn 9-13% with automated leveraged strategies.",
         font_size=14, color=WHITE, bold=True)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "THE SOLUTION: SUPERLOOP", font_size=32, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.5),
         "The institutional vault infrastructure layer for automated, compliant, leveraged RWA strategies.",
         font_size=16, color=LIGHT_GREY)

features = [
    ("One-Click Looping",
     "Flash-loan powered single-transaction looping. Deposit stablecoins, vault auto-loops to target leverage. No manual intervention."),
    ("Modular Architecture",
     "Plug-and-play modules for Aave Horizon, Centrifuge, Morpho, DEX swaps. Add new protocols without rewriting core logic."),
    ("Built-In Compliance",
     "KYC gating at vault entry. Transfer-restricted shares. Compatible with ERC-1404, ERC-3643, and Chainlink ACE."),
    ("ERC-4626 Standard",
     "Fully composable tokenized vault. Async deposit/withdraw queues. Performance fee via share dilution. Institutional-grade accounting."),
]

for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 5.8)
    top = Inches(1.8 + row * 2.1)
    card = add_rounded_rect(slide, left, top, Inches(5.5), Inches(1.8), DARK_CARD, ACCENT)
    add_text(slide, left + Inches(0.25), top + Inches(0.15), Inches(5), Inches(0.4),
             title, font_size=18, color=ACCENT, bold=True)
    add_text(slide, left + Inches(0.25), top + Inches(0.65), Inches(5), Inches(1),
             desc, font_size=13, color=LIGHT_GREY)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "MARKET OPPORTUNITY", font_size=32, color=ACCENT, bold=True)

kpis = [
    ("Tokenized RWAs On-Chain", "$35B+", "Growing 300%+ YoY"),
    ("Aave Horizon TVL", "$1B", "Reached Feb 2026"),
    ("Morpho Looping Volume", "64%", "Of total $3B loans"),
    ("BlackRock BUIDL AUM", "$2.3B+", "Largest tokenized fund"),
]
for i, (title, value, sub) in enumerate(kpis):
    add_kpi_card(slide, Inches(0.8 + i * 3.05), Inches(1.3), Inches(2.8), Inches(1.3),
                 title, value, sub)

add_text(slide, Inches(0.8), Inches(3.1), Inches(11), Inches(0.5),
         "Key Trend: RWA issuers (Superstate, Centrifuge, BlackRock/Securitize, VanEck) are deploying tokens, but lack DeFi-native yield infrastructure.",
         font_size=14, color=WHITE, bold=True)

partners_text = [
    ("RWA Issuers: Superstate (USTB, USCC) | Centrifuge (JTRSY, JAAA) | Securitize (BUIDL) | VanEck (VBILL) | Hashnote (USYC)", LIGHT_GREY),
    ("Lending Markets: Aave Horizon ($1B) | Morpho ($3B) | Spark ($3.5B)", LIGHT_GREY),
    ("Stablecoin Borrowing: USDC (Circle) | GHO (Aave DAO) | RLUSD (Ripple)", LIGHT_GREY),
    ("Infrastructure: Chainlink (NAVLink + ACE) | Chaos Labs (Risk Oracles) | LlamaRisk (Due Diligence)", LIGHT_GREY),
]
add_multiline(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(3), partners_text, font_size=14, bullet=True)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (NON-TECHNICAL)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "HOW IT WORKS", font_size=32, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         "From User Deposit to Leveraged RWA Yield in One Transaction", font_size=16, color=LIGHT_GREY)

steps = [
    ("1", "DEPOSIT", "User deposits USDC\ninto Superloop Vault"),
    ("2", "ACQUIRE", "Vault swaps USDC for\nRWA token (e.g. USTB)"),
    ("3", "SUPPLY", "Supply RWA token as\ncollateral on Horizon"),
    ("4", "BORROW", "Borrow USDC against\nRWA collateral"),
    ("5", "LOOP", "Repeat steps 2-4\nusing flash loans"),
    ("6", "EARN", "Vault earns leveraged\nspread (9-13% APR)"),
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.5 + i * 2.05)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.55), Inches(1.8), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    card = add_rounded_rect(slide, left, Inches(2.5), Inches(1.85), Inches(2.2), DARK_CARD, MED_GREY)
    add_text(slide, left + Inches(0.1), Inches(2.6), Inches(1.65), Inches(0.35),
             title, font_size=13, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), Inches(3.05), Inches(1.65), Inches(1.5),
             desc, font_size=11, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # Arrow between steps (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(1.9), Inches(3.3),
                                        Inches(0.2), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

add_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.8),
         "The entire loop executes atomically in a single transaction using flash loans.\nNo manual intervention. No multi-step risk. Gas efficient.",
         font_size=14, color=WHITE, bold=True)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — YIELD ECONOMICS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "YIELD ECONOMICS", font_size=32, color=ACCENT, bold=True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Leveraged APR = L x Y - (L-1) x R      where L = leverage, Y = RWA yield, R = borrow rate",
         font_size=16, color=HIGHLIGHT, bold=True)

# Table header
table_data = [
    ["Leverage", "5% Yield / 3% Borrow", "5% Yield / 2% Borrow", "8% Yield / 3% Borrow"],
    ["1x (base)", "5.0%", "5.0%", "8.0%"],
    ["2x", "7.0%", "8.0%", "13.0%"],
    ["3x", "9.0%", "11.0%", "18.0%"],
    ["4x", "11.0%", "14.0%", "23.0%"],
    ["5x (max)", "13.0%", "17.0%", "28.0%"],
]

tbl_shape = slide.shapes.add_table(len(table_data), len(table_data[0]),
                                     Inches(0.8), Inches(1.7), Inches(10), Inches(3))
tbl = tbl_shape.table

for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        if r == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_CARD if r % 2 == 1 else RGBColor(0x15, 0x17, 0x25)

highlights = [
    ("Conservative target: 3x leverage on Treasury collateral = 9-11% APR", ACCENT2, True),
    ("Flash loan source: Morpho (0% fee) or Balancer (0% fee) to maximize thin RWA spreads", LIGHT_GREY, False),
    ("Revenue model: 10-20% performance fee + 0.5-2% management fee on AUM", LIGHT_GREY, False),
]
add_multiline(slide, Inches(0.8), Inches(5.1), Inches(11), Inches(1.5), highlights, font_size=14, bullet=True)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE OVERVIEW
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "ARCHITECTURE OVERVIEW", font_size=32, color=ACCENT, bold=True)

# Architecture layers as stacked cards
layers = [
    ("USER LAYER", "ERC-4626 Vault | Async Deposit Queue | Withdraw Manager | Share Accounting", ACCENT),
    ("STRATEGY LAYER", "VaultOperator calls operate() with ModuleExecutionData[] | Flash Loan Wrapping | Callback Execution", ACCENT2),
    ("MODULE LAYER", "AaveV3Supply | AaveV3Borrow | AaveV3Flashloan | MorphoFlashloan | UniversalDex | VaultSupply | Emode", HIGHLIGHT),
    ("ACCOUNTING LAYER", "AaveV3AccountantPlugin | NAV Oracle Integration | Performance Fee Calculation | Position Tracking", RGBColor(0xE0, 0x60, 0xFF)),
    ("PROTOCOL LAYER", "Aave Horizon | Aave V3 | Morpho | Centrifuge | Chainlink NAVLink | DEX Aggregators", RGBColor(0xFF, 0x60, 0x60)),
]

for i, (title, desc, color) in enumerate(layers):
    top = Inches(1.2 + i * 1.1)
    card = add_rounded_rect(slide, Inches(0.8), top, Inches(11.5), Inches(0.9), DARK_CARD, color)
    add_text(slide, Inches(1.1), top + Inches(0.05), Inches(2.5), Inches(0.35),
             title, font_size=14, color=color, bold=True)
    add_text(slide, Inches(1.1), top + Inches(0.4), Inches(11), Inches(0.4),
             desc, font_size=12, color=LIGHT_GREY)

add_text(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
         "Key insight: Each module is a standalone contract. Adding Horizon = deploying existing AaveV3 modules with a new PoolAddressesProvider. Zero code changes.",
         font_size=13, color=WHITE, bold=True)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — EXISTING INFRASTRUCTURE (VERIFIED)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "WHAT'S ALREADY BUILT (VERIFIED)", font_size=32, color=ACCENT, bold=True)

# Two columns
left_items = [
    ("AaveV3SupplyModule", "pool.supply(asset, amount, address(this), 0)", True),
    ("AaveV3BorrowModule", "pool.borrow(asset, amount, INTEREST_RATE_MODE, 0, address(this))", True),
    ("AaveV3FlashloanModule", "pool.flashLoanSimple() with callback execution data", True),
    ("AaveV3RepayModule", "Repay borrowed positions", True),
    ("AaveV3WithdrawModule", "Withdraw supplied collateral", True),
    ("AaveV3EmodeModule", "Set efficiency mode for correlated pairs", True),
    ("MorphoFlashloanModule", "morpho.flashLoan() - zero fee flash loans", True),
]

right_items = [
    ("UniversalDexModule", "Multi-DEX swap execution (1inch, Uniswap, etc.)", True),
    ("AaveV3AccountantPlugin", "Tracks lend + borrow positions via getUserReserveData", True),
    ("SuperloopVault (ERC-4626)", "Deposit, withdraw, share accounting, performance fees", True),
    ("DepositManager", "Async deposit queue with request/resolve pattern", True),
    ("WithdrawManager", "Multi-queue: Instant, Priority, Deferred, General", True),
    ("Callback System", "Flash loan callbacks route through operateSelf()", True),
    ("ModuleRegistry", "Whitelist-based module registration and validation", True),
]

for i, (title, desc, _) in enumerate(left_items):
    top = Inches(1.2 + i * 0.75)
    add_text(slide, Inches(0.8), top, Inches(5.5), Inches(0.3),
             title, font_size=13, color=ACCENT2, bold=True)
    add_text(slide, Inches(0.8), top + Inches(0.28), Inches(5.5), Inches(0.35),
             desc, font_size=11, color=LIGHT_GREY)

for i, (title, desc, _) in enumerate(right_items):
    top = Inches(1.2 + i * 0.75)
    add_text(slide, Inches(6.8), top, Inches(5.5), Inches(0.3),
             title, font_size=13, color=ACCENT2, bold=True)
    add_text(slide, Inches(6.8), top + Inches(0.28), Inches(5.5), Inches(0.35),
             desc, font_size=11, color=LIGHT_GREY)

add_rounded_rect(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.55), DARK_CARD, ACCENT)
add_text(slide, Inches(1.1), Inches(6.65), Inches(11), Inches(0.45),
         "All AaveV3 modules take poolAddressesProvider in constructor. Horizon integration = new deployments with Horizon provider address. ZERO code changes needed.",
         font_size=13, color=ACCENT, bold=True)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — AAVE HORIZON INTEGRATION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "AAVE HORIZON INTEGRATION", font_size=32, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Horizon = Aave V3.3 fork for permissioned RWA lending. $1B TVL. Same interface, different PoolAddressesProvider.",
         font_size=14, color=LIGHT_GREY)

# Dual structure diagram
add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.3), Inches(2.5), DARK_CARD, ACCENT)
add_text(slide, Inches(1.0), Inches(1.7), Inches(4.5), Inches(0.35),
         "PERMISSIONED SIDE (Collateral)", font_size=16, color=ACCENT, bold=True)
perm_items = [
    "USTB - Superstate US Treasuries",
    "JTRSY - Centrifuge / Janus Henderson Treasuries",
    "JAAA - Centrifuge AAA CLOs",
    "USYC - Hashnote / Circle Yield Fund",
    "VBILL - VanEck Treasury Fund",
    "KYC required per issuer whitelist",
    "aTokens are NON-TRANSFERABLE",
]
add_multiline(slide, Inches(1.0), Inches(2.1), Inches(4.8), Inches(1.8),
              [(item, LIGHT_GREY) for item in perm_items], font_size=11, bullet=True)

add_rounded_rect(slide, Inches(6.5), Inches(1.6), Inches(5.3), Inches(2.5), DARK_CARD, ACCENT2)
add_text(slide, Inches(6.7), Inches(1.7), Inches(4.5), Inches(0.35),
         "PERMISSIONLESS SIDE (Lending)", font_size=16, color=ACCENT2, bold=True)
open_items = [
    "USDC - Circle (primary liquidity)",
    "GHO - Aave DAO stablecoin",
    "RLUSD - Ripple institutional stablecoin",
    "Anyone can supply stablecoins to earn yield",
    "Yield comes from institutional borrowers",
    "No KYC needed for stablecoin supply",
    "Stablecoins CANNOT be used as collateral",
]
add_multiline(slide, Inches(6.7), Inches(2.1), Inches(4.8), Inches(1.8),
              [(item, LIGHT_GREY) for item in open_items], font_size=11, bullet=True)

# Superloop integration
add_rounded_rect(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.5), DARK_CARD, HIGHLIGHT)
add_text(slide, Inches(1.0), Inches(4.5), Inches(10), Inches(0.35),
         "SUPERLOOP VAULT FLOW ON HORIZON", font_size=16, color=HIGHLIGHT, bold=True)

flow_text = (
    "1. Deploy AaveV3SupplyModule(HORIZON_POOL_ADDRESSES_PROVIDER)\n"
    "2. Deploy AaveV3BorrowModule(HORIZON_POOL_ADDRESSES_PROVIDER)\n"
    "3. Vault address gets KYC-whitelisted by RWA issuer (e.g., Superstate for USTB)\n"
    "4. Operator calls operate(): FlashLoan USDC -> Swap to USTB -> Supply USTB to Horizon -> Borrow USDC -> Repay flash loan\n"
    "5. AaveV3AccountantPlugin tracks position: lend (USTB aTokens) - borrow (USDC debt) = net vault value\n"
    "6. Vault share price reflects leveraged Treasury yield (9-13% target APR)"
)
add_text(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(1.8),
         flow_text, font_size=12, color=LIGHT_GREY)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — CENTRIFUGE INTEGRATION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "CENTRIFUGE INTEGRATION", font_size=32, color=ACCENT, bold=True)

add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Access tokenized structured credit: AAA CLOs, US Treasuries, diversified credit. ERC-7540 async vault standard.",
         font_size=14, color=LIGHT_GREY)

# Centrifuge products
products = [
    ("JTRSY", "Janus Henderson\nUS Treasury\nExposure", "~4.5-5% APR\nSenior tranche\nAAA-equivalent"),
    ("JAAA", "Janus Henderson\nAAA CLOs\n(loan obligations)", "~5-7% APR\nFirst-claim priority\n$1B+ allocated"),
    ("ACRDX", "Apollo\nDiversified Credit\n(via Grove)", "~7-10% APR\nInstitutional credit\nPrivate markets"),
]

for i, (ticker, desc, yield_info) in enumerate(products):
    left = Inches(0.8 + i * 4.0)
    card = add_rounded_rect(slide, left, Inches(1.6), Inches(3.7), Inches(2.0), DARK_CARD, ACCENT)
    add_text(slide, left + Inches(0.2), Inches(1.7), Inches(3.3), Inches(0.35),
             ticker, font_size=22, color=ACCENT, bold=True)
    add_text(slide, left + Inches(0.2), Inches(2.1), Inches(1.5), Inches(1.3),
             desc, font_size=12, color=LIGHT_GREY)
    add_text(slide, left + Inches(1.8), Inches(2.1), Inches(1.7), Inches(1.3),
             yield_info, font_size=12, color=ACCENT2)

# Integration requirements
add_text(slide, Inches(0.8), Inches(3.9), Inches(10), Inches(0.35),
         "NEW MODULE REQUIRED: CentrifugeSupplyModule", font_size=16, color=HIGHLIGHT, bold=True)

centrifuge_flow = [
    ("ERC-7540 Async Flow: requestDeposit() -> wait for epoch (24h) -> claimDeposit() -> receive tranche tokens", LIGHT_GREY),
    ("Maps to Superloop's existing DepositManager/WithdrawManager async pattern", LIGHT_GREY),
    ("Tranche tokens can then be supplied to Aave Horizon as collateral for leveraged yield", LIGHT_GREY),
    ("Grove Finance routes Sky's $1B allocation through Centrifuge V3 into JAAA and ACRDX", LIGHT_GREY),
]
add_multiline(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(2.5), centrifuge_flow,
              font_size=13, bullet=True)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — COMPLIANCE FRAMEWORK
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "COMPLIANCE FRAMEWORK", font_size=32, color=ACCENT, bold=True)

# Three layers
layer_data = [
    ("LAYER 1: Entity-Level", ACCENT, [
        "SPV / Segregated Portfolio Company (Cayman Islands preferred)",
        "Section 3(c)(7) fund exemption for qualified purchasers",
        "Vault smart contract address KYC'd by each RWA issuer",
        "Institutional-grade custody (Anchorage, BitGo, Fireblocks)",
    ]),
    ("LAYER 2: Investor-Level", ACCENT2, [
        "KYC/AML verification via Sumsub or Chainlink ACE (CCID)",
        "Accredited / Qualified Purchaser attestation",
        "Whitelist enforcement at DepositManager entry point",
        "Transfer-restricted vault shares (onlyPrivileged modifier)",
    ]),
    ("LAYER 3: On-Chain Enforcement", HIGHLIGHT, [
        "ERC-1404: detectTransferRestriction() check before transfers",
        "ERC-3643 (T-REX): Vault address registered in Identity Registry",
        "Chainlink ACE: Cross-chain identity (CCID) + policy verification",
        "Non-transferable aTokens in Horizon maintain permissioning integrity",
    ]),
]

for i, (title, color, items) in enumerate(layer_data):
    left = Inches(0.8 + i * 4.0)
    card = add_rounded_rect(slide, left, Inches(1.2), Inches(3.7), Inches(3.5), DARK_CARD, color)
    add_text(slide, left + Inches(0.2), Inches(1.3), Inches(3.3), Inches(0.35),
             title, font_size=15, color=color, bold=True)
    add_multiline(slide, left + Inches(0.2), Inches(1.75), Inches(3.3), Inches(2.8),
                  [(item, LIGHT_GREY) for item in items], font_size=11, bullet=True)

add_rounded_rect(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.5), DARK_CARD, ACCENT)
add_text(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(0.35),
         "SUPERLOOP'S EXISTING COMPLIANCE INFRASTRUCTURE", font_size=14, color=ACCENT, bold=True)
existing = [
    ("privilegedAddresses mapping: Already restricts share transfers to whitelisted addresses (SuperloopVault.sol:330)", LIGHT_GREY),
    ("onlyPrivileged modifier: Controls transfer() and transferFrom() on vault shares", LIGHT_GREY),
    ("DepositManager: Natural entry point for KYC gating (add onlyCompliant check to requestDeposit)", LIGHT_GREY),
]
add_multiline(slide, Inches(1.0), Inches(5.6), Inches(11), Inches(1), existing, font_size=11, bullet=True)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — EXACT INTEGRATIONS NEEDED
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "EXACT INTEGRATIONS REQUIRED AT SCALE", font_size=32, color=ACCENT, bold=True)

integrations = [
    ("ZERO CODE CHANGES (Deploy Only)", ACCENT2, [
        "AaveV3SupplyModule  ->  new instance with HORIZON_POOL_ADDRESSES_PROVIDER",
        "AaveV3BorrowModule  ->  new instance with HORIZON_POOL_ADDRESSES_PROVIDER",
        "AaveV3RepayModule   ->  new instance with HORIZON_POOL_ADDRESSES_PROVIDER",
        "AaveV3WithdrawModule ->  new instance with HORIZON_POOL_ADDRESSES_PROVIDER",
        "AaveV3EmodeModule   ->  new instance with HORIZON_POOL_ADDRESSES_PROVIDER",
        "AaveV3FlashloanModule ->  new instance (if Horizon supports flashLoanSimple for stablecoins)",
        "AaveV3AccountantPlugin -> new instance with Horizon PoolAddressesProvider + RWA lend/borrow assets",
    ]),
    ("NEW MODULES NEEDED", HIGHLIGHT, [
        "CentrifugeSupplyModule  ->  ERC-7540 requestDeposit/claimDeposit/requestRedeem",
        "RWAComplianceGate       ->  KYC verification, investor tier checks, Chainlink ACE integration",
        "NAVOracleAdapter        ->  Chainlink NAVLink price feed integration for RWA pricing",
        "CentrifugeAccountantPlugin -> Track pending/claimable positions + epoch-based NAV",
        "GroveRouterModule       ->  Capital routing through Grove/Sky allocation network",
    ]),
    ("INFRASTRUCTURE & OPS", RGBColor(0xE0, 0x60, 0xFF), [
        "Keeper/Operator Bot     ->  Monitor rates, auto-deleverage, epoch claims, health factor tracking",
        "Legal Entity (SPV)      ->  Required for each RWA issuer whitelisting",
        "KYC Provider Integration ->  Sumsub, Coinbase Verifications, or Chainlink CCID",
        "Risk Parameter Engine   ->  Chaos Labs / LlamaRisk integration for dynamic LTV/rate adjustment",
    ]),
]

y_offset = Inches(1.15)
for title, color, items in integrations:
    add_text(slide, Inches(0.8), y_offset, Inches(10), Inches(0.35),
             title, font_size=14, color=color, bold=True)
    y_offset += Inches(0.35)
    for item in items:
        add_text(slide, Inches(1.1), y_offset, Inches(11), Inches(0.28),
                 item, font_size=10.5, color=LIGHT_GREY)
        y_offset += Inches(0.26)
    y_offset += Inches(0.15)

add_slide_number(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — EXECUTION FLOW (TECHNICAL)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "EXECUTION FLOW: SINGLE-TX RWA LOOP", font_size=32, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Based on proven EthenaLoopTest pattern in production codebase", font_size=14, color=LIGHT_GREY)

code = """
// Operator calls: superloop.operate(finalExecutionData)

Step 1: resolveDepositRequests(USDC, depositAmount, callback=[
  Step 2: morphoFlashLoan(USTB, flashAmount, callback=[
    Step 3: aaveV3Supply(USTB, type(uint256).max)     // Supply all USTB to Horizon
    Step 4: aaveV3Borrow(USDC, borrowAmount)          // Borrow USDC against USTB collateral
    Step 5: dexSwap(USDC -> USTB, borrowAmount)       // Swap borrowed USDC to more USTB
    Step 6: aaveV3Supply(USTB, type(uint256).max)     // Re-supply new USTB (loop iteration)
    Step 7: aaveV3Borrow(USDC, repayAmount)           // Borrow to repay flash loan
    // Flash loan callback auto-approves repayment
  ])
])

// Result: Vault holds 4-5x leveraged USTB position on Horizon
// AccountantPlugin reads: aToken balance (USTB) - variableDebt (USDC) = net value
"""

card = add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.2), RGBColor(0x12, 0x14, 0x20), ACCENT)
add_text(slide, Inches(1.1), Inches(1.6), Inches(11), Inches(4),
         code.strip(), font_size=12, color=ACCENT2)

add_text(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.8),
         "Atomic execution: all steps succeed or all revert. No partial state. No multi-block risk.\nProven pattern: EthenaLoopTest (test/core/integration/EthenaLoop.t.sol) demonstrates this exact flow.",
         font_size=13, color=WHITE, bold=True)
add_slide_number(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — RISK MANAGEMENT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "RISK MANAGEMENT", font_size=32, color=ACCENT, bold=True)

risks = [
    ("Liquidation Risk", "Low", "Conservative LTV (60-70%). Treasuries have near-zero volatility.\nAuto-deleverage triggers at health factor 1.5."),
    ("Borrow Rate Spike", "Medium", "Dynamic monitoring. Auto-deleverage when spread < 50bps.\nMultiple stablecoin sources (USDC, GHO, RLUSD)."),
    ("Oracle Staleness", "Low", "Chainlink NAVLink with price bounds validation.\nDON rejects out-of-band NAV updates."),
    ("Smart Contract", "Low", "Battle-tested Aave V3.3 (audited + formally verified).\nSuperloop modules are minimal surface area."),
    ("Redemption Delay", "Medium", "Cash reserve buffer (configurable BPS).\nMulti-queue WithdrawManager: instant, priority, deferred."),
    ("Regulatory", "Medium", "SPV legal structure. Jurisdictional optionality.\nCompliance at token + vault + investor levels."),
]

for i, (risk, severity, mitigation) in enumerate(risks):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.0)
    top = Inches(1.1 + row * 1.85)
    sev_color = ACCENT2 if severity == "Low" else HIGHLIGHT
    card = add_rounded_rect(slide, left, top, Inches(5.7), Inches(1.6), DARK_CARD, sev_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), Inches(4), Inches(0.3),
             risk, font_size=14, color=WHITE, bold=True)
    add_text(slide, left + Inches(4.2), top + Inches(0.1), Inches(1.2), Inches(0.3),
             severity, font_size=12, color=sev_color, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text(slide, left + Inches(0.2), top + Inches(0.5), Inches(5.3), Inches(1),
             mitigation, font_size=11, color=LIGHT_GREY)

add_slide_number(slide, 14, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — SUPPORTED ASSETS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "SUPPORTED ASSETS & PARTNERS", font_size=32, color=ACCENT, bold=True)

# RWA Collateral table
rwa_header = ["Asset", "Issuer", "Underlying", "Platform", "Yield"]
rwa_data = [
    ["USTB", "Superstate", "US Govt Securities", "Aave Horizon", "~4.5%"],
    ["USCC", "Superstate", "Crypto Carry Fund", "Aave Horizon", "~5-8%"],
    ["JTRSY", "Centrifuge", "Janus Henderson Treasuries", "Aave Horizon", "~4.5%"],
    ["JAAA", "Centrifuge", "Janus Henderson AAA CLOs", "Aave Horizon", "~5-7%"],
    ["USYC", "Hashnote/Circle", "Int'l Short Duration Yield", "Aave Horizon", "~4.5%"],
    ["VBILL", "VanEck/Securitize", "VanEck Treasury Fund", "Aave Horizon", "~4.5%"],
    ["BUIDL", "BlackRock/Securitize", "Treasury Money Market", "Pipeline", "~4.5%"],
]
all_data = [rwa_header] + rwa_data
tbl_shape = slide.shapes.add_table(len(all_data), len(rwa_header),
                                     Inches(0.8), Inches(1.1), Inches(11.5), Inches(3.8))
tbl = tbl_shape.table

for r, row in enumerate(all_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.name = "Calibri"
        if r == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            p.font.color.rgb = WHITE if val != "Pipeline" else HIGHLIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_CARD if r % 2 == 1 else RGBColor(0x15, 0x17, 0x25)

stables = [
    ("Borrowable Stablecoins: USDC (Circle) | GHO (Aave DAO, ~$312M cap) | RLUSD (Ripple)", LIGHT_GREY),
    ("Pipeline Partners: Franklin Templeton | WisdomTree | Ethena | OpenEden | Ant Digital", LIGHT_GREY),
]
add_multiline(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(1.5), stables, font_size=14, bullet=True)
add_slide_number(slide, 15, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 16 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "COMPETITIVE POSITIONING", font_size=32, color=ACCENT, bold=True)

comp_header = ["Competitor", "What They Do", "Superloop Advantage"]
comp_data = [
    ["Gauntlet", "Curated levered RWA vaults on Morpho", "Gauntlet is a vault manager, not infra.\nSuperloop is the vault rails they could build on."],
    ["Huma Finance", "Defensive Looping vaults (PayFi)", "Narrow focus (PST tokens only).\nSuperloop is multi-RWA, multi-protocol."],
    ["Sommelier", "Dynamic strategy vaults", "Cosmos bridge dependency.\nSuperloop is native EVM, no bridge risk."],
    ["Morpho Vaults", "Permissionless lending markets", "Morpho is a lending primitive.\nSuperloop composes Morpho + Aave + Centrifuge."],
    ["Summer.fi", "Frontend for leveraged positions", "Frontend only, no vault abstraction.\nNo compliance. No managed strategies."],
]

all_comp = [comp_header] + comp_data
tbl_shape = slide.shapes.add_table(len(all_comp), 3, Inches(0.8), Inches(1.2), Inches(11.5), Inches(4))
tbl = tbl_shape.table

# Set column widths
tbl.columns[0].width = Inches(2.3)
tbl.columns[1].width = Inches(4.0)
tbl.columns[2].width = Inches(5.2)

for r, row in enumerate(all_comp):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.name = "Calibri"
        if r == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_CARD if r % 2 == 1 else RGBColor(0x15, 0x17, 0x25)

add_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(1.2),
         "Superloop's moat: Not another lending market or frontend. We're the institutional vault infrastructure layer that sits ABOVE Aave Horizon, Centrifuge, Morpho, and Grove -- composing them into managed, compliant, leveraged RWA strategies.",
         font_size=14, color=ACCENT, bold=True)
add_slide_number(slide, 16, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 17 — REVENUE MODEL
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "REVENUE MODEL", font_size=32, color=ACCENT, bold=True)

rev_streams = [
    ("Performance Fee", "10-20% of yield", "Share dilution on profit (already implemented in SuperloopVault)", ACCENT),
    ("Management Fee", "0.5-2% annually", "AUM-based via exchange rate adjustment in accountant", ACCENT2),
    ("Instant Withdraw Fee", "0.1-0.5%", "Premium for skipping withdrawal queue (already implemented)", HIGHLIGHT),
]

for i, (name, rate, desc, color) in enumerate(rev_streams):
    left = Inches(0.8 + i * 4.0)
    card = add_rounded_rect(slide, left, Inches(1.2), Inches(3.7), Inches(1.8), DARK_CARD, color)
    add_text(slide, left + Inches(0.2), Inches(1.3), Inches(3.3), Inches(0.3),
             name, font_size=16, color=color, bold=True)
    add_text(slide, left + Inches(0.2), Inches(1.7), Inches(3.3), Inches(0.35),
             rate, font_size=22, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.2), Inches(2.2), Inches(3.3), Inches(0.6),
             desc, font_size=11, color=LIGHT_GREY)

# Example economics
add_text(slide, Inches(0.8), Inches(3.4), Inches(10), Inches(0.4),
         "EXAMPLE: $100M TVL, 10% Leveraged APR, 15% Performance Fee, 1% Management Fee",
         font_size=16, color=WHITE, bold=True)

econ_kpis = [
    ("Gross Yield", "$10M / year", ""),
    ("Performance Fee", "$1.5M / year", "15% of yield"),
    ("Management Fee", "$1.0M / year", "1% of AUM"),
    ("Total Protocol Revenue", "$2.5M / year", "Per $100M TVL"),
]
for i, (title, value, sub) in enumerate(econ_kpis):
    add_kpi_card(slide, Inches(0.8 + i * 3.05), Inches(4.1), Inches(2.8), Inches(1.3),
                 title, value, sub)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
         "Scaling path: $500M TVL = ~$12.5M annual revenue. $1B TVL = ~$25M annual revenue.\nRevenue scales linearly with TVL. Multiple vault strategies multiply the opportunity.",
         font_size=14, color=LIGHT_GREY)
add_slide_number(slide, 17, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 18 — ROADMAP
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "ROADMAP", font_size=32, color=ACCENT, bold=True)

phases = [
    ("PHASE 1: Horizon Loop Vault", "Weeks 1-6", ACCENT, [
        "Deploy Aave V3 modules with Horizon PoolAddressesProvider",
        "Build + test RWA loop execution (adapt EthenaLoopTest)",
        "Integrate Chainlink NAVLink oracle adapter",
        "Legal entity setup + KYC vault with Superstate (USTB)",
        "RWAComplianceGate contract for investor verification",
        "Security audit + testnet deployment",
    ]),
    ("PHASE 2: Centrifuge Integration", "Weeks 7-12", ACCENT2, [
        "Build CentrifugeSupplyModule (ERC-7540 async)",
        "Keeper bot for epoch monitoring + auto-claim",
        "Combined strategy: Centrifuge tranche tokens -> Horizon -> leverage",
        "KYC vault with Centrifuge (JAAA, JTRSY issuers)",
        "Multi-RWA vault support (diversified positions)",
        "Mainnet launch with institutional partners",
    ]),
    ("PHASE 3: Scale & Expand", "Weeks 13-20", HIGHLIGHT, [
        "Grove partnership for Sky $1B capital allocation",
        "Multi-RWA diversified vaults (USTB + JAAA + ACRDX)",
        "Senior/Junior tranche vaults for risk stratification",
        "Cross-chain RWA vaults (Base, Arbitrum) via Centrifuge V3",
        "Chainlink ACE (CCID) for cross-chain compliance",
        "Morpho RWA market integration for additional lending venues",
    ]),
]

for i, (title, timeline, color, items) in enumerate(phases):
    left = Inches(0.8 + i * 4.0)
    card = add_rounded_rect(slide, left, Inches(1.1), Inches(3.7), Inches(5.5), DARK_CARD, color)
    add_text(slide, left + Inches(0.2), Inches(1.2), Inches(3.3), Inches(0.3),
             title, font_size=14, color=color, bold=True)
    add_text(slide, left + Inches(0.2), Inches(1.55), Inches(3.3), Inches(0.25),
             timeline, font_size=12, color=MED_GREY, bold=True)
    add_multiline(slide, left + Inches(0.2), Inches(1.9), Inches(3.3), Inches(4.5),
                  [(item, LIGHT_GREY) for item in items], font_size=11, bullet=True, spacing=1.4)

add_slide_number(slide, 18, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 19 — WHY WHITELIST SUPERLOOP (FOR RWA ISSUERS)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "WHY WHITELIST SUPERLOOP", font_size=32, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
         "For RWA Issuers: What Superloop Brings to Your Token", font_size=16, color=LIGHT_GREY)

reasons = [
    ("Demand Generation",
     "Superloop vaults create SUSTAINED demand for your RWA tokens. Leveraged strategies consume 3-5x more tokens per dollar of user deposits. At $100M vault TVL, we buy $300-500M of your tokens.",
     ACCENT),
    ("Non-Custodial & Transparent",
     "All positions are on-chain, auditable in real-time. Smart contracts execute deterministically. We cannot move funds. Aave V3 is battle-tested with $30B+ historical TVL.",
     ACCENT2),
    ("Compliance-First Architecture",
     "Our vault contract completes your KYC/AML requirements. We enforce investor qualification at deposit entry. Vault shares are transfer-restricted. Compatible with ERC-1404, ERC-3643, Chainlink ACE.",
     HIGHLIGHT),
    ("Institutional Distribution",
     "Superloop serves as a distribution channel for your RWA tokens to qualified purchasers. Single KYC relationship (vault entity) unlocks access for all verified vault depositors.",
     RGBColor(0xE0, 0x60, 0xFF)),
]

for i, (title, desc, color) in enumerate(reasons):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.0)
    top = Inches(1.5 + row * 2.3)
    card = add_rounded_rect(slide, left, top, Inches(5.7), Inches(2.0), DARK_CARD, color)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), Inches(5.3), Inches(0.35),
             title, font_size=16, color=color, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.5), Inches(5.3), Inches(1.4),
             desc, font_size=12, color=LIGHT_GREY)

add_text(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
         "Ask: Whitelist the Superloop vault contract address as a qualified purchaser for your RWA token.",
         font_size=16, color=ACCENT, bold=True)
add_slide_number(slide, 19, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# SLIDE 20 — CONTACT / CLOSING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "SUPERLOOP", font_size=56, color=ACCENT, bold=True)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
         "Building the Institutional RWA Yield Layer", font_size=28, color=WHITE)

add_rect(slide, Inches(1), Inches(3.5), Inches(4), Inches(0.03), ACCENT)

summary_points = [
    ("70%+ of infrastructure already built and production-tested", WHITE, True),
    ("Aave Horizon integration requires zero code changes -- deploy only", WHITE, True),
    ("$35B+ addressable RWA market growing 300%+ YoY", WHITE, True),
    ("9-13% APR target on Treasury-backed collateral at 3-4x leverage", WHITE, True),
    ("Compliance-first: KYC gating, transfer restrictions, institutional legal structure", WHITE, True),
]
add_multiline(slide, Inches(1), Inches(4.0), Inches(11), Inches(2.5), summary_points, font_size=16, spacing=1.5, bullet=True)

add_text(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
         "Superlend  |  superloop-core-contracts  |  February 2026  |  Confidential",
         font_size=14, color=MED_GREY, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 20, TOTAL_SLIDES)

# ── Save ──────────────────────────────────────────────────────
output_dir = "/home/user/superloop-core-contracts/docs"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "Superloop_RWA_Strategy_Deck.pptx")
prs.save(output_path)
print(f"Deck saved to: {output_path}")
